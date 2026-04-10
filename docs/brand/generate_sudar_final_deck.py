from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "docs" / "brand" / "Sudar_Final_Deck.pptx"
LOGO = ROOT / "assets" / "logos" / "Sudar Logo_transparent blue.png"


COLOR_BG = RGBColor(13, 16, 38)  # #0D1026
COLOR_BG_LIGHT = RGBColor(247, 248, 252)  # #F7F8FC
COLOR_PRIMARY = RGBColor(47, 42, 138)  # #2F2A8A
COLOR_PRIMARY_SOFT = RGBColor(94, 90, 215)  # #5E5AD7
COLOR_ACCENT = RGBColor(255, 122, 69)  # #FF7A45
COLOR_HIGHLIGHT = RGBColor(255, 209, 102)  # #FFD166
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_DARK = RGBColor(26, 32, 44)
COLOR_MUTED = RGBColor(95, 105, 130)


def add_bg(slide, dark=False):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG if dark else COLOR_BG_LIGHT


def add_title(slide, title, subtitle=None, dark=False):
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE if dark else COLOR_PRIMARY
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11), Inches(0.8))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.size = Pt(20)
        srun.font.color.rgb = COLOR_HIGHLIGHT if dark else COLOR_MUTED


def add_bullets(slide, items, x=0.95, y=2.3, w=11.2, h=4.5, dark=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(10)
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_WHITE if dark else COLOR_DARK


def add_logo(slide, dark=False):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(10.3), Inches(0.35), width=Inches(2.5))
    else:
        # Fallback mark if asset path changes.
        box = slide.shapes.add_textbox(Inches(10.2), Inches(0.5), Inches(2.5), Inches(0.6))
        tf = box.text_frame
        tf.text = "Sudar"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE if dark else COLOR_PRIMARY


def add_footer(slide, text, dark=False):
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(6.8), Inches(12.0), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY_SOFT if dark else RGBColor(220, 224, 235)
    line.line.fill.background()

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(10), Inches(0.4))
    tf = foot.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = COLOR_HIGHLIGHT if dark else COLOR_MUTED


def content_slide(prs, title, bullets, subtitle=None, dark=False, footer="Sudar | Final Brand Deck"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, dark=dark)
    add_logo(slide, dark=dark)
    add_title(slide, title, subtitle=subtitle, dark=dark)
    add_bullets(slide, bullets, dark=dark)
    add_footer(slide, footer, dark=dark)


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, dark=True)
    add_logo(slide, dark=True)
    add_title(
        slide,
        "Sudar",
        subtitle="The Operating System for Learning",
        dark=True,
    )
    add_bullets(
        slide,
        [
            "Learns with you, for you.",
            "AI-native learning platform with longitudinal learner memory",
            "Final integrated brand and product narrative",
        ],
        y=2.7,
        h=2.4,
        dark=True,
    )
    add_footer(slide, "April 2026 | Founder/Product Vision Deck", dark=True)

    # 2. Why now
    content_slide(
        prs,
        "Why Sudar, Why Now",
        [
            "Most LMS platforms track completion, but do not remember learners.",
            "Teams need measurable outcomes without large production budgets.",
            "Sudar closes the gap with adaptive delivery + memory-aware tutoring.",
        ],
        subtitle="From one-size-fits-all training to persistent personalization",
    )

    # 3. Brand core
    content_slide(
        prs,
        "Brand Core",
        [
            "Mission: Democratize high-quality personalized learning.",
            "Promise: Learning that remembers each learner and adapts in real time.",
            "Strategic position: Speed + Evidence + Memory-aware personalization.",
            "Voice: Human, clear, outcome-first, claim-safe.",
        ],
    )

    # 4. Logo narrative
    content_slide(
        prs,
        "Logo Narrative",
        [
            "“=” represents equal access to high-quality learning.",
            "Dynamic motion forms the “S” for learning flow.",
            "Center AI star/flame represents adaptive intelligence.",
            "Visual expression: premium-minimal with warm-human cues.",
        ],
        subtitle="A symbol that maps directly to Sudar's product philosophy",
    )

    # 5. Product system
    content_slide(
        prs,
        "Three Surfaces, One Intelligence Layer",
        [
            "Sudar Studio: rapid authoring, path orchestration, admin workflows.",
            "Sudar Learn: learner-first experience with modality flexibility.",
            "Sudar Intelligence: adaptive engine, next-best-action, tutor memory.",
            "Shared Supabase layer unifies profiles, events, and personalization.",
        ],
    )

    # 6. Differentiators
    content_slide(
        prs,
        "What Makes Sudar Different",
        [
            "Digital Learner Twin stores longitudinal behavior and preferences.",
            "Sudar tutor is reactive and proactive, with cross-session memory.",
            "Author once, deliver across modalities (text, flashcards, video, audio, etc.).",
            "ALP direction supports augmentation of existing LMS ecosystems.",
        ],
    )

    # 7. Audience fit
    content_slide(
        prs,
        "Who This Is For (Next 90 Days)",
        [
            "L&D Managers: launch effective training quickly with lean teams.",
            "HR/Talent Leaders: connect learning programs to workforce outcomes.",
            "Shared message: Build faster, personalize by default, prove progress.",
        ],
    )

    # 8. Proof stack
    content_slide(
        prs,
        "Proof Stack and Claim Discipline",
        [
            "Validated: shipped architecture and documented capabilities.",
            "Evidence-informed: design grounded in adaptive learning research.",
            "Future-work: pilots and roadmap claims clearly marked as planned.",
            "Outcome: ambitious narrative with high trust and low hype risk.",
        ],
        dark=True,
    )

    # 9. Visual tokens
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, dark=False)
    add_logo(slide, dark=False)
    add_title(slide, "Visual System at a Glance", "Option A palette + typography", dark=False)
    add_footer(slide, "Design tokens: docs/brand/design-tokens-v1.md", dark=False)

    swatches = [
        ("Indigo Base", COLOR_PRIMARY),
        ("Indigo Support", COLOR_PRIMARY_SOFT),
        ("Ember Accent", COLOR_ACCENT),
        ("Warm Highlight", COLOR_HIGHLIGHT),
        ("Deep Night", COLOR_BG),
        ("Soft Cloud", COLOR_BG_LIGHT),
    ]
    x0, y0, w, h = 0.95, 2.3, 1.8, 1.1
    for i, (label, color) in enumerate(swatches):
        x = x0 + (i % 3) * 2.15
        y = y0 + (i // 3) * 2.05
        rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.color.rgb = RGBColor(220, 224, 235)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y + 1.2), Inches(w + 0.2), Inches(0.6))
        tf = tb.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = COLOR_DARK

    type_box = slide.shapes.add_textbox(Inches(7.6), Inches(2.35), Inches(4.9), Inches(2.8))
    ttf = type_box.text_frame
    ttf.clear()
    p1 = ttf.paragraphs[0]
    p1.text = "Typography"
    p1.font.bold = True
    p1.font.size = Pt(26)
    p1.font.color.rgb = COLOR_PRIMARY
    for line in [
        "Headlines: Manrope",
        "UI / Body: Inter",
        "4/8 spacing rhythm, rounded surfaces,",
        "light + dark theme parity required",
    ]:
        p = ttf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(4)

    # 10. Final vision
    content_slide(
        prs,
        "Final Product Vision",
        [
            "Sudar becomes the adaptive learning operating system for modern organizations.",
            "The platform remembers each learner, not just their completions.",
            "Teams can author once, personalize continuously, and scale outcomes responsibly.",
            "Brand, product, and research are now unified under one clear system.",
        ],
        subtitle="Build faster. Teach better. Learn personally.",
        dark=True,
        footer="Sudar | Learns with you, for you.",
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    make_deck()
